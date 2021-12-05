import lottery_festival as lot
from pptx import Presentation
from pptx.util import Cm

lottery = lot.lotteryFestival()

command = 'y'
count = 0
lottery.display_all_attendee()

while(1):
    in_str = input('当選者を何人にしますか?: ')
    # 入力が整数じゃなければやり直し
    try:
        winner_count = int(in_str)
    except ValueError:
        print('エラーメッセージ: 数字を入力してください')
        continue
    
    count += 1
    print('-----抽選', count, '回目', '-----', sep="")
    lottery.select_winner(winner_count)
    lottery.display_all_winner()
    #lottery.display_all_attendee()
    
    command = input('続けて抽選を行いますか？y/n: ')
    if command != 'y':
        print("抽選を終了します")
        break
    
    
# 表紙スライドを追加
def add_front_cover(ppt):
    slide_0_layout = ppt.slide_layouts[0]
    slide_0 = ppt.slides.add_slide(slide_0_layout)
    # タイトルの縦横を設定
    slide_title = slide_0.placeholders[0]
    slide_title.width = Cm(25.4)
    slide_title.height = Cm(6.63)
    # タイトルのパラグラフを追加
    pg = slide_title.text_frame.paragraphs[0]
    # 文字を追加
    pg.text = '当選者発表'
    # フォントサイズを変更
    pg.font.size = Cm(3)
    

# 当選者スライドを追加
def add_all_winner(ppt):
    #スライドを追加
    for i in range(len(lottery.list_all_winner_number)):
        #追加するスライドのレイアウトを選択
        slide_layout = ppt.slide_layouts[5]
        #スライドを追加
        slide = ppt.slides.add_slide(slide_layout)
        # プレースホルダー0(タイトル)を追加
        slide_title = slide.placeholders[0]
        # タイトルの縦横を設定
        slide_title.width = Cm(25.4)
        slide_title.height = Cm(6.63)
        # タイトルのテキストフレームを追加
        text_frame = slide_title.text_frame
        # テキストフレームのパラグラフを追加
        pg = text_frame.paragraphs[0]
        # 文字を追加
        pg.text = str(lottery.list_all_winner_number[0])
        # フォントサイズを変更
        pg.font.size = Cm(5)

    
#Presentaionオブジェクトを生成
ppt = Presentation()

#スライドのサイズを指定
ppt.slide_width = Cm(33.867)
ppt.slide_height = Cm(19.05)

add_front_cover(ppt)
add_all_winner(ppt)

# ファイルを閉じていない場合例外を出す
try:
    ppt.save('./result_lottery.pptx')
except PermissionError:
    print('powerpointファイルを閉じてください')
else:
    print('powerpointファイルを作成しました')