from pptx import Presentation
from pptx.util import Cm
from pptx.enum.text import MSO_AUTO_SIZE   # MOS_AUTO_SIZEクラスのインポート

class powerpoint():
    def __init__(self):
        #Presentaionオブジェクトを生成
        self.ppt = Presentation()
        #スライドのサイズを指定
        self.ppt.slide_width = Cm(33.867)
        self.ppt.slide_height = Cm(19.05)
        
    # 表紙スライドを追加
    def add_front_cover(self, title):
        slide_0_layout = self.ppt.slide_layouts[0]
        slide_0 = self.ppt.slides.add_slide(slide_0_layout)
        # タイトルの縦横を設定
        slide_title = slide_0.placeholders[0]
        
        # スライドの縦横を設定
        slide_title.width = Cm(25.4)
        slide_title.height = Cm(6.63)
        
        base = slide_title.height
        slide_title.height = int(self.ppt.slide_height * 0.7)
        slide_title.width = int(slide_title.width * (slide_title.height / base))
        # 左右中央揃え
        slide_title.left = int((self.ppt.slide_width - slide_title.width) / 2)
        # 下揃え
        slide_title.top = int((self.ppt.slide_height / 1.3 - slide_title.height))
        
        # タイトルのパラグラフを追加
        pg = slide_title.text_frame.paragraphs[0]
        # 文字を追加
        pg.text = title
        # フォントサイズを変更
        pg.font.size = Cm(4)
        self.ppt.slide_height
        
    
    # 当選者スライドを追加 パターン1
    def add_each_list_1(self, list_text):
        #スライドを追加
        for i in range(len(list_text)):
            #追加するスライドのレイアウトを選択
            slide_layout = self.ppt.slide_layouts[5]
            #スライドを追加
            slide = self.ppt.slides.add_slide(slide_layout)
            # プレースホルダー0(タイトル)を追加
            slide_title = slide.placeholders[0]
            
            # スライドの縦横を設定
            slide_title.width = Cm(25.4)
            slide_title.height = Cm(6.63)
            
            base = Cm(19.05)
            slide_title.height = int(self.ppt.slide_height * 0.7)
            slide_title.width = int(slide_title.width * (slide_title.height / base))
            # 左右中央揃え
            slide_title.left = int((self.ppt.slide_width - slide_title.width) / 2)
            # 下揃え
            slide_title.top = int((self.ppt.slide_height / 1.3 - slide_title.height))
            
            # タイトルのテキストフレームを追加
            text_frame = slide_title.text_frame
            # TextFrameのテキストサイズを自動調整
            text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
            # テキストフレームのパラグラフを追加
            pg = text_frame.paragraphs[0]
            # 文字を追加
            pg.text = str(list_text[i][0]) + '\n' + str(list_text[i][1])
            #フォントサイズを変更
            pg.font.size = Cm(3)
            
    
    # 当選者スライドを追加 パターン1
    def add_each_list_2(self, list_text):
        slide_layout = self.ppt.slide_layouts[5]
        slide = self.ppt.slides.add_slide(slide_layout)
        slide_title = slide.placeholders[0]
        
        # スライドの縦横を設定
        slide_title.width = Cm(25.4)
        slide_title.height = Cm(6.63)
        
        base = Cm(19.05)
        slide_title.height = int(self.ppt.slide_height * 0.7)
        slide_title.width = int(slide_title.width * (slide_title.height / base))
        # 左右中央揃え
        slide_title.left = int((self.ppt.slide_width - slide_title.width) / 2)
        # 下揃え
        slide_title.top = int((self.ppt.slide_height / 1.3 - slide_title.height))
        
        pg = slide_title.text_frame.paragraphs[0]
        pg.text = ''
        for text in list_text:
            pg.text += str(text) + '\n'
        

    def save_pptx(self, title):
        # ファイルを閉じていない場合例外を出す
        try:
            self.ppt.save('./' + title + '.pptx')
        except PermissionError:
            print('エラーメッセージ: ' + title + '.pptx' + 'を閉じてください')
        else:
            print(title + '.pptx' + 'を作成しました')
